import os
import io
import pandas as pd
from pdfminer.high_level import extract_text_to_fp
from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
from dotenv import load_dotenv

# Load Gemini LLM client
try:
    import google.generativeai as genai
except ImportError:
    raise ImportError("google-generativeai package not installed. Run: pip install google-generativeai")

# Optional: PowerPoint support
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    print("Warning: python-pptx not found. Install it if you need PPT support.")
    Presentation = None
    PPTX_AVAILABLE = False

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
LLM_MODEL = "gemini-2.5-flash"

# Flask app (fix: template_folder="templates")
app = Flask(__name__, template_folder="templates")  
CORS(app)

# Initialize Gemini client
client = None
if GEMINI_API_KEY:
    try:
        client = genai.Client(api_key=GEMINI_API_KEY)
        print("Gemini client initialized successfully.")
    except Exception as e:
        print(f"Error initializing Gemini client: {e}")
else:
    print("CRITICAL: GEMINI_API_KEY not found in .env. LLM will not work.")

# --- File parsing ---
def extract_text_from_file(file):
    filename = file.filename.lower()
    file_stream = io.BytesIO(file.read())

    try:
        if filename.endswith(".pdf"):
            output_string = io.StringIO()
            extract_text_to_fp(file_stream, output_string)
            return output_string.getvalue()

        elif filename.endswith(".csv"):
            df = pd.read_csv(file_stream)
            return df.to_string(index=False)

        elif filename.endswith((".xlsx", ".xls")):
            df = pd.read_excel(file_stream)
            return df.to_string(index=False)

        elif filename.endswith((".pptx", ".ppt")):
            if not PPTX_AVAILABLE:
                return None
            prs = Presentation(file_stream)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip() != "":
                        text_content.append(shape.text.strip())
            return "\n".join(text_content)

        else:
            return None

    except Exception as e:
        print(f"Error parsing {filename}: {e}")
        return None

# --- LLM call ---
def run_gemini_generation(system_prompt, user_prompt):
    if not client:
        return "LLM client not initialized. Check GEMINI_API_KEY on server."
    try:
        response = client.models.generate_content(
            model=LLM_MODEL,
            contents=[user_prompt],
            config={"system_instruction": system_prompt}
        )
        return response.text
    except Exception as e:
        print(f"LLM Error: {e}")
        return f"LLM Error: {e}"

# --- API endpoint ---
@app.route("/api/rag/ask", methods=["POST"])
def ask_question():
    if "document" not in request.files:
        return jsonify({"answer": "No document uploaded."}), 400

    file = request.files["document"]
    query = request.form.get("query", "").strip()

    if file.filename == "" or query == "":
        return jsonify({"answer": "File and query are required."}), 400

    doc_content = extract_text_from_file(file)
    if not doc_content:
        if file.filename.lower().endswith((".pptx", ".ppt")) and not PPTX_AVAILABLE:
            return jsonify({"answer": "python-pptx library missing for PowerPoint files."}), 400
        return jsonify({"answer": "Could not extract text from file (Unsupported format or corrupted file)."}), 400

    # --- RAG prompt ---
    system_prompt = (
        "You are an expert RAG system. Answer the user's QUESTION strictly based "
        "on the provided CONTEXT. If the information is not in the CONTEXT, "
        "state clearly that it is not available in the document. Respond concisely."
    )
    user_prompt = f"CONTEXT:\n---\n{doc_content}\n---\n\nQUESTION: {query}\n\nANSWER:"

    answer = run_gemini_generation(system_prompt, user_prompt)

    if answer.startswith("LLM Error:"):
        return jsonify({"answer": answer}), 500

    return jsonify({"answer": answer})

# --- Serve frontend ---
@app.route("/")
def home():
    return render_template("index.html")  # Flask now looks inside templates/

# --- Run server ---
if __name__ == "__main__":
    print("Starting Flask server at http://127.0.0.1:5000")
    app.run(debug=True)
